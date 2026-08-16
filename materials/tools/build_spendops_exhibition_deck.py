from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUTPUT = ROOT / "materials" / "deliverables" / "SpendOps_Dashboard_展示資料.pptx"
ICON = ROOT / "materials" / "images" / "icon" / "icon.png"
TECH_DOC_LINK = "SpendOps_Dashboard_技術解説.docx"

W = 13.333
H = 7.5
FONT = "Noto Sans JP"

DARK = "101713"
LIME = "DFFF78"
ACCENT = "B8DF4B"
PALE = "EFF8CF"
OFFWHITE = "F6F7F4"
WHITE = "FFFFFF"
MUTED = "68736C"
LINE = "D9E0DA"
RED = "B52D38"
BLUE = "1D63A5"
AMBER = "C47A16"


def rgb(hex_color: str) -> RGBColor:
    return RGBColor.from_string(hex_color)


def set_east_asian_font(run, name: str = FONT) -> None:
    run.font.name = name
    rpr = run._r.get_or_add_rPr()
    for child in list(rpr):
        if child.tag.endswith("}ea"):
            rpr.remove(child)
    east_asian = OxmlElement("a:ea")
    east_asian.set("typeface", name)
    rpr.append(east_asian)


def fill_shape(shape, color: str, transparency: int = 0) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)
    shape.fill.transparency = transparency


def no_line(shape) -> None:
    shape.line.fill.background()


def box(slide, x, y, w, h, color, radius=True, line_color=None, line_width=1):
    geometry = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(geometry, Inches(x), Inches(y), Inches(w), Inches(h))
    fill_shape(shape, color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width)
    else:
        no_line(shape)
    return shape


def circle(slide, x, y, d, color, line_color=None, line_width=1):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    fill_shape(shape, color)
    if line_color:
        shape.line.color.rgb = rgb(line_color)
        shape.line.width = Pt(line_width)
    else:
        no_line(shape)
    return shape


def line(slide, x1, y1, x2, y2, color, width=1.5, arrow=False):
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    connector.line.color.rgb = rgb(color)
    connector.line.width = Pt(width)
    if arrow:
        connector.line.end_arrowhead = True
    return connector


def text(
    slide,
    value,
    x,
    y,
    w,
    h,
    *,
    size=20,
    color=DARK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin=0.02,
    line_spacing=1.0,
    link=None,
    name=None,
):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    if name:
        shape.name = name
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(margin)
    tf.margin_bottom = Inches(margin)
    tf.vertical_anchor = valign
    paragraphs = str(value).split("\n")
    for index, paragraph_value in enumerate(paragraphs):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = paragraph_value
        set_east_asian_font(run)
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
        if link:
            run.hyperlink.address = link
    return shape


def rich_text(slide, runs, x, y, w, h, *, size=20, color=DARK, valign=MSO_ANCHOR.TOP, align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = p.space_after = Pt(0)
    for spec in runs:
        run = p.add_run()
        run.text = spec[0]
        set_east_asian_font(run)
        run.font.size = Pt(spec[1] if len(spec) > 1 and spec[1] else size)
        run.font.bold = bool(spec[2]) if len(spec) > 2 else False
        run.font.color.rgb = rgb(spec[3] if len(spec) > 3 and spec[3] else color)
        if len(spec) > 4 and spec[4]:
            run.hyperlink.address = spec[4]
    return shape


def bullet_list(slide, items, x, y, w, h, *, size=17, color=DARK, bullet_color=ACCENT, gap=0.52):
    for i, item in enumerate(items):
        yy = y + i * gap
        circle(slide, x, yy + 0.13, 0.12, bullet_color)
        text(slide, item, x + 0.24, yy, w - 0.24, gap, size=size, color=color, valign=MSO_ANCHOR.MIDDLE)


def pill(slide, label, x, y, w, *, fill=LIME, color=DARK, size=12):
    shape = box(slide, x, y, w, 0.36, fill, True)
    text(slide, label, x, y, w, 0.36, size=size, color=color, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    return shape


def title_block(slide, eyebrow, title_value, subtitle=None, *, dark=False):
    fg = WHITE if dark else DARK
    pill(slide, eyebrow, 0.72, 0.44, max(1.2, min(3.3, len(eyebrow) * 0.17 + 0.55)), fill=LIME, color=DARK)
    text(slide, title_value, 0.72, 0.98, 11.9, 0.82, size=31, color=fg, bold=True, valign=MSO_ANCHOR.MIDDLE)
    if subtitle:
        text(slide, subtitle, 0.74, 1.74, 11.65, 0.52, size=15, color=("D7DED9" if dark else MUTED), valign=MSO_ANCHOR.MIDDLE)


def footer(slide, number, minutes, *, dark=False):
    color = "A8B2AC" if dark else MUTED
    line(slide, 0.72, 7.08, 12.62, 7.08, "344039" if dark else LINE, 0.8)
    text(slide, "SpendOps Dashboard", 0.72, 7.12, 3.2, 0.2, size=9, color=color)
    text(slide, f"目安 {minutes}分", 9.65, 7.12, 1.8, 0.2, size=9, color=color, align=PP_ALIGN.RIGHT)
    text(slide, f"{number:02d} / 12", 11.5, 7.12, 1.12, 0.2, size=9, color=color, align=PP_ALIGN.RIGHT)


def add_slide(prs, bg):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background.fill
    background.solid()
    background.fore_color.rgb = rgb(bg)
    return slide


def add_demo_badge(slide, x, y, w=1.25):
    pill(slide, "表示例（デモ）", x, y, w, fill=PALE, color=DARK, size=10)


def slide_01(prs):
    s = add_slide(prs, DARK)
    pill(s, "展示用｜約20分｜自由閲覧", 0.72, 0.52, 2.65, fill=LIME, color=DARK, size=11)
    text(s, "SpendOps\nDashboard", 0.74, 1.38, 7.5, 1.7, size=43, color=WHITE, bold=True)
    text(s, "支払いを、1つの画面で振り返る", 0.76, 3.18, 7.35, 0.7, size=27, color=LIME, bold=True)
    text(s, "CSVを読み込む → 月ごとに比較 → 次の行動を決める", 0.77, 4.03, 7.25, 0.55, size=16, color="D7DED9")
    box(s, 8.82, 0.75, 3.75, 5.7, LIME, True)
    circle(s, 9.35, 1.4, 2.7, WHITE)
    if ICON.exists():
        s.shapes.add_picture(str(ICON), Inches(9.52), Inches(1.57), width=Inches(2.36), height=Inches(2.36))
    text(s, "明細を集める", 9.1, 4.45, 3.15, 0.38, size=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "↓", 10.28, 4.85, 0.8, 0.42, size=21, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    text(s, "気づきを行動へ", 9.1, 5.3, 3.15, 0.42, size=17, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    box(s, 0.74, 5.48, 7.4, 0.72, "243027", True)
    text(s, "現在はAWS基盤を停止中。実装内容と設計判断を資料で展示します。", 0.98, 5.62, 6.92, 0.35, size=13, color="D7DED9", valign=MSO_ANCHOR.MIDDLE)
    footer(s, 1, 1, dark=True)


def slide_02(prs):
    s = add_slide(prs, OFFWHITE)
    title_block(s, "01｜課題", "「今月どうだった？」が、支払い先ごとに分かれている", "明細の形式や締め日が違うため、全体像をつかむまでに手間がかかります。")
    sources = [("PayPay", 0.78, 2.65), ("JCB", 0.78, 3.55), ("三井住友VISA", 0.78, 4.45)]
    for label, x, y in sources:
        box(s, x, y, 2.6, 0.62, WHITE, True, LINE)
        text(s, label, x, y, 2.6, 0.62, size=18, color=DARK, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
        line(s, x + 2.6, y + 0.31, 4.44, 3.86, "AAB4AE", 1.2)
    circle(s, 4.36, 3.12, 1.5, DARK)
    text(s, "ばらばら", 4.36, 3.52, 1.5, 0.42, size=17, color=LIME, bold=True, align=PP_ALIGN.CENTER)
    line(s, 5.95, 3.86, 6.9, 3.86, DARK, 2.2, arrow=True)
    box(s, 7.03, 2.58, 5.35, 3.15, PALE, True)
    text(s, "必要なのは\n“同じ見方”", 7.45, 2.95, 4.5, 1.25, size=29, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    bullet_list(s, ["合計を1つにする", "前月や目安と比べる", "使いすぎた分野を見つける"], 7.78, 4.45, 3.9, 1.35, size=14, gap=0.42)
    footer(s, 2, 1.5)


def slide_03(prs):
    s = add_slide(prs, DARK)
    title_block(s, "02｜解決", "3つのCSVを、1つの月次レポートへ", "対応対象：PayPay、JCB、三井住友VISA。銀行CSVは今回の対象外です。", dark=True)
    x = 0.85
    for i, label in enumerate(["PayPay.csv", "JCB.csv", "VISA.csv"]):
        box(s, x, 2.73 + i * 0.92, 2.55, 0.64, "243027", True, "48564D")
        text(s, label, x, 2.73 + i * 0.92, 2.55, 0.64, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    line(s, 3.55, 3.95, 4.55, 3.95, LIME, 2.6, arrow=True)
    circle(s, 4.63, 3.25, 1.45, LIME)
    text(s, "共通形式\nへ変換", 4.63, 3.56, 1.45, 0.72, size=15, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    line(s, 6.18, 3.95, 7.0, 3.95, LIME, 2.6, arrow=True)
    box(s, 7.14, 2.4, 5.35, 3.82, WHITE, True)
    add_demo_badge(s, 10.93, 2.62, 1.2)
    text(s, "2026年6月", 7.52, 2.76, 2.7, 0.35, size=13, color=MUTED)
    text(s, "支出合計", 7.52, 3.24, 2.2, 0.32, size=12, color=MUTED)
    rich_text(s, [("86,420", 32, True, DARK), ("円", 16, True, DARK)], 7.5, 3.52, 3.05, 0.55)
    pill(s, "前月比 -8.4%", 10.5, 3.53, 1.55, fill=PALE, color=BLUE, size=11)
    text(s, "カテゴリ上位", 7.52, 4.38, 2.0, 0.3, size=12, color=MUTED)
    bars = [("食費", 0.82, "24,800円"), ("日用品", 0.54, "16,300円"), ("交通", 0.38, "11,600円")]
    for i, (label, rate, amount) in enumerate(bars):
        yy = 4.78 + i * 0.38
        text(s, label, 7.52, yy, 0.72, 0.26, size=11, color=DARK)
        box(s, 8.35, yy + 0.04, 2.55, 0.13, "E2E8E3", False)
        box(s, 8.35, yy + 0.04, 2.55 * rate, 0.13, ACCENT, False)
        text(s, amount, 11.03, yy - 0.01, 1.05, 0.27, size=10, color=DARK, align=PP_ALIGN.RIGHT)
    footer(s, 3, 2, dark=True)


def slide_04(prs):
    s = add_slide(prs, OFFWHITE)
    title_block(s, "03｜使い方", "使い方は3ステップ", "読み込みから分析まではブラウザ内で完結。ログイン後は必要な集計だけを保存できます。")
    steps = [
        ("01", "CSVを用意", "各サービスから利用明細CSVを出力します。", "書き換え不要"),
        ("02", "ファイルを選ぶ", "画面でCSVを選ぶと、形式を自動判定します。", "原本は送信しない"),
        ("03", "レポートを見る", "月別・年間・全期間を切り替えて確認します。", "比較と次の行動"),
    ]
    for i, (num, heading, body, tag) in enumerate(steps):
        x = 0.78 + i * 4.15
        box(s, x, 2.52, 3.75, 3.35, WHITE, True, LINE)
        circle(s, x + 0.24, 2.77, 0.62, DARK)
        text(s, num, x + 0.24, 2.89, 0.62, 0.25, size=13, color=LIME, bold=True, align=PP_ALIGN.CENTER)
        text(s, heading, x + 0.28, 3.58, 3.1, 0.5, size=23, color=DARK, bold=True)
        text(s, body, x + 0.28, 4.25, 3.05, 0.78, size=15, color=MUTED)
        pill(s, tag, x + 0.28, 5.25, 2.05, fill=PALE, color=DARK, size=11)
        if i < 2:
            line(s, x + 3.78, 4.2, x + 4.08, 4.2, ACCENT, 2.4, arrow=True)
    box(s, 0.78, 6.14, 12.0, 0.64, DARK, True)
    rich_text(s, [("ポイント：", 13, True, LIME), ("「アップロード」と表示しても、CSV原本は端末内で解析されます。", 13, False, WHITE)], 1.02, 6.31, 11.55, 0.28)
    footer(s, 4, 1.5)


def slide_05(prs):
    s = add_slide(prs, OFFWHITE)
    title_block(s, "04｜読み方", "最初に見るのは、合計・比較・内訳", "月別を中心に、必要に応じて年間・全期間へ広げます。")
    box(s, 0.78, 2.45, 8.35, 4.08, WHITE, True, LINE)
    pill(s, "月別", 1.05, 2.73, 0.92, fill=DARK, color=LIME, size=11)
    pill(s, "年間", 2.04, 2.73, 0.92, fill=PALE, color=DARK, size=11)
    pill(s, "全期間", 3.03, 2.73, 1.05, fill=PALE, color=DARK, size=11)
    add_demo_badge(s, 7.57, 2.73, 1.18)
    text(s, "支出合計", 1.08, 3.34, 1.3, 0.32, size=12, color=MUTED)
    rich_text(s, [("86,420", 30, True, DARK), ("円", 15, True, DARK)], 1.06, 3.65, 2.6, 0.52)
    text(s, "比較", 3.48, 3.34, 0.9, 0.32, size=12, color=MUTED)
    text(s, "前月より 7,880円少ない", 3.47, 3.69, 2.8, 0.42, size=16, color=BLUE, bold=True)
    text(s, "カテゴリ内訳", 1.08, 4.52, 1.5, 0.3, size=12, color=MUTED)
    data = [("食費", 0.86, ACCENT), ("日用品", 0.6, "75A456"), ("交通", 0.42, "A7B8AF"), ("その他", 0.3, "D8DFDA")]
    for i, (label, rate, color) in enumerate(data):
        yy = 4.92 + i * 0.34
        text(s, label, 1.08, yy - 0.03, 0.76, 0.25, size=10, color=DARK)
        box(s, 1.94, yy, 5.55, 0.12, "E9EDE9", False)
        box(s, 1.94, yy, 5.55 * rate, 0.12, color, False)
    text(s, "→ 明細まで確認し、必要ならカテゴリを修正", 5.2, 5.94, 3.55, 0.28, size=11, color=MUTED, align=PP_ALIGN.RIGHT)
    bullet_list(s, ["合計：今月の規模", "比較：増減とその理由", "内訳：行動を変える候補"], 9.72, 2.75, 2.65, 2.0, size=16, gap=0.67)
    box(s, 9.52, 5.1, 2.85, 1.12, PALE, True)
    text(s, "見る順番を固定すると\n短時間でも迷いません", 9.76, 5.31, 2.37, 0.7, size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    footer(s, 5, 2)


def slide_06(prs):
    s = add_slide(prs, PALE)
    title_block(s, "05｜気づき", "比較は、平均を見るためではなく行動を決めるため", "差額・変化・集中しているカテゴリから、次に試せる行動を示します。")
    text(s, "今月", 0.92, 2.6, 1.0, 0.3, size=13, color=MUTED)
    rich_text(s, [("86,420", 38, True, DARK), ("円", 17, True, DARK)], 0.9, 2.92, 3.2, 0.68)
    text(s, "前月", 4.2, 2.6, 1.0, 0.3, size=13, color=MUTED)
    rich_text(s, [("94,300", 38, True, DARK), ("円", 17, True, DARK)], 4.18, 2.92, 3.2, 0.68)
    line(s, 0.94, 3.92, 7.55, 3.92, DARK, 2)
    box(s, 0.94, 3.8, 5.72, 0.24, ACCENT, False)
    pill(s, "-8.4%", 6.9, 3.73, 0.95, fill=DARK, color=LIME, size=12)
    box(s, 8.42, 2.55, 3.85, 3.45, DARK, True)
    text(s, "次に試すこと", 8.78, 2.9, 3.15, 0.42, size=16, color=LIME, bold=True)
    text(s, "外食が前月より\n5,600円増えています", 8.78, 3.52, 3.05, 0.82, size=21, color=WHITE, bold=True)
    line(s, 8.78, 4.6, 11.85, 4.6, "47534B", 1)
    text(s, "来週は外食を1回だけ\n置き換えて変化を確認", 8.78, 4.86, 3.05, 0.78, size=16, color="D7DED9")
    text(s, "※ 表示例。提案は支出データに基づく一般的な振り返りで、金融助言ではありません。", 0.94, 6.32, 7.0, 0.28, size=10, color=MUTED)
    footer(s, 6, 1.5)


def slide_07(prs):
    s = add_slide(prs, DARK)
    title_block(s, "06｜学習", "カテゴリ修正は、次回の手間を減らす", "本人が直した分類を、同じ利用先の次回取込へ反映します。", dark=True)
    text(s, "取込直後", 0.88, 2.48, 1.4, 0.3, size=12, color="A8B2AC")
    box(s, 0.86, 2.85, 4.2, 1.1, "243027", True, "48564D")
    text(s, "コンビニABC", 1.12, 3.1, 1.85, 0.4, size=17, color=WHITE, bold=True)
    pill(s, "その他", 3.72, 3.18, 0.95, fill="49564D", color=WHITE, size=11)
    line(s, 5.35, 3.4, 6.44, 3.4, LIME, 2.5, arrow=True)
    circle(s, 6.55, 2.84, 1.12, LIME)
    text(s, "本人が\n修正", 6.55, 3.08, 1.12, 0.6, size=14, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    line(s, 7.78, 3.4, 8.63, 3.4, LIME, 2.5, arrow=True)
    text(s, "次回から", 8.72, 2.48, 1.4, 0.3, size=12, color="A8B2AC")
    box(s, 8.7, 2.85, 3.76, 1.1, WHITE, True)
    text(s, "コンビニABC", 8.98, 3.1, 1.85, 0.4, size=17, color=DARK, bold=True)
    pill(s, "食費", 11.12, 3.18, 0.9, fill=PALE, color=DARK, size=11)
    box(s, 0.88, 4.75, 11.55, 1.2, "1A241E", True)
    bullet_list(s, ["ルールは本人単位で保存", "利用先は正規化して照合", "他ユーザーの分類には影響しない"], 1.18, 5.0, 10.9, 0.9, size=14, color=WHITE, bullet_color=LIME, gap=0.34)
    footer(s, 7, 1.5, dark=True)


def slide_08(prs):
    s = add_slide(prs, OFFWHITE)
    title_block(s, "07｜個人情報保護", "CSV原本を送らず、必要な情報だけを保存する", "安全性の中心は「何を守るか」だけでなく、「そもそも持たないか」です。")
    box(s, 0.78, 2.5, 5.73, 3.85, DARK, True)
    text(s, "端末内に留まる", 1.12, 2.88, 4.9, 0.48, size=21, color=LIME, bold=True)
    text(s, "クラウドへ送らない", 1.14, 3.4, 4.8, 0.3, size=12, color="A8B2AC")
    bullet_list(s, ["CSV原本・ファイル名", "元の行データ", "カード番号・口座情報", "パスワード・認証コード"], 1.15, 3.92, 4.75, 1.75, size=15, color=WHITE, bullet_color=LIME, gap=0.46)
    box(s, 6.82, 2.5, 5.73, 3.85, WHITE, True, LINE)
    text(s, "認証後に保存", 7.17, 2.88, 4.95, 0.48, size=21, color=DARK, bold=True)
    text(s, "画面復元と集計に必要な最小項目", 7.19, 3.4, 4.85, 0.3, size=12, color=MUTED)
    bullet_list(s, ["本人ID・日付・金額", "正規化した利用先・カテゴリ", "取得元・取込ID", "月別集計・分類ルールのハッシュ"], 7.2, 3.92, 4.8, 1.75, size=15, color=DARK, bullet_color=ACCENT, gap=0.46)
    text(s, "詳細なデータ項目・暗号化・認証方式は別紙の技術資料で確認できます。", 0.8, 6.57, 11.74, 0.3, size=11, color=MUTED, align=PP_ALIGN.CENTER)
    footer(s, 8, 2)


def slide_09(prs):
    s = add_slide(prs, DARK)
    title_block(s, "08｜匿名比較", "実ユーザー比較は「他5人以上・完全月」だけ", "条件を満たさないときは、実データを混ぜず「参考例」と明示した合成データへ切り替えます。", dark=True)
    box(s, 0.9, 2.63, 3.55, 2.95, "243027", True, "48564D")
    text(s, "比較候補", 1.22, 2.95, 2.9, 0.36, size=14, color="A8B2AC", bold=True)
    for i in range(6):
        circle(s, 1.32 + (i % 3) * 0.85, 3.48 + (i // 3) * 0.86, 0.55, LIME if i < 5 else "49564D")
    text(s, "本人は除外", 1.15, 5.2, 3.0, 0.25, size=11, color="A8B2AC", align=PP_ALIGN.CENTER)
    line(s, 4.66, 4.1, 5.62, 4.1, LIME, 2.4, arrow=True)
    circle(s, 5.76, 3.43, 1.36, LIME)
    text(s, "判定", 5.76, 3.83, 1.36, 0.32, size=16, color=DARK, bold=True, align=PP_ALIGN.CENTER)
    line(s, 7.25, 4.1, 8.18, 4.1, LIME, 2.4, arrow=True)
    box(s, 8.34, 2.61, 4.05, 1.32, WHITE, True)
    text(s, "他5人以上 + 完全月", 8.68, 2.9, 3.35, 0.36, size=17, color=DARK, bold=True)
    pill(s, "匿名集計を表示", 9.2, 3.39, 2.25, fill=PALE, color=DARK, size=11)
    box(s, 8.34, 4.25, 4.05, 1.32, "243027", True, "48564D")
    text(s, "人数不足 / 一部期間", 8.68, 4.54, 3.35, 0.36, size=17, color=WHITE, bold=True)
    pill(s, "参考例（合成）", 9.2, 5.03, 2.25, fill="49564D", color=WHITE, size=11)
    box(s, 0.9, 6.05, 11.49, 0.58, "1A241E", True)
    text(s, "登録時に匿名集計への利用同意を取得し、比較結果から個人を表示しません。", 1.16, 6.2, 10.95, 0.26, size=12, color="D7DED9", align=PP_ALIGN.CENTER)
    footer(s, 9, 2, dark=True)


def slide_10(prs):
    s = add_slide(prs, OFFWHITE)
    title_block(s, "09｜検証", "守りは設計だけでなく、70件の自動テストで確認", "フロントエンド46件とLambda 24件。正常系・異常系・境界条件を継続確認します。")
    text(s, "70", 0.88, 2.56, 2.65, 1.2, size=67, color=DARK, bold=True)
    text(s, "自動テスト成功", 0.95, 3.7, 2.55, 0.4, size=16, color=MUTED, bold=True)
    text(s, "フロントエンド", 4.0, 2.72, 2.15, 0.3, size=13, color=MUTED)
    box(s, 4.0, 3.13, 5.65, 0.42, "E1E6E2", True)
    box(s, 4.0, 3.13, 5.65 * 46 / 70, 0.42, ACCENT, True)
    text(s, "46件", 9.82, 3.1, 0.8, 0.35, size=14, color=DARK, bold=True)
    text(s, "Lambda", 4.0, 3.88, 2.15, 0.3, size=13, color=MUTED)
    box(s, 4.0, 4.29, 5.65, 0.42, "E1E6E2", True)
    box(s, 4.0, 4.29, 5.65 * 24 / 70, 0.42, "75A456", True)
    text(s, "24件", 9.82, 4.26, 0.8, 0.35, size=14, color=DARK, bold=True)
    box(s, 0.9, 5.05, 11.47, 1.25, PALE, True)
    text(s, "確認している例", 1.22, 5.36, 1.4, 0.28, size=12, color=MUTED, bold=True)
    rich_text(s, [("CSV形式差・文字コード", 13, True, DARK), ("　/　", 13, False, MUTED), ("重複取込", 13, True, DARK), ("　/　", 13, False, MUTED), ("JWT・本人分離", 13, True, DARK), ("　/　", 13, False, MUTED), ("比較の5人条件", 13, True, DARK), ("　/　", 13, False, MUTED), ("分類学習", 13, True, DARK)], 2.8, 5.28, 9.1, 0.55, valign=MSO_ANCHOR.MIDDLE)
    text(s, "※ テスト件数は2026年7月23日時点。すべての脅威や端末環境を保証するものではありません。", 0.94, 6.47, 11.2, 0.28, size=10, color=MUTED)
    footer(s, 10, 1.5)


def slide_11(prs):
    s = add_slide(prs, PALE)
    title_block(s, "10｜現在地", "実装済み。ただし、現在はクラウド機能を停止中", "提出後にAWS基盤を削除したため、公開サイト・ログイン・クラウド保存・実ユーザー比較は動作しません。")
    box(s, 0.78, 2.55, 5.7, 3.7, WHITE, True)
    pill(s, "実装済み", 1.1, 2.88, 1.25, fill=PALE, color=DARK, size=12)
    bullet_list(s, ["3種CSVの解析・正規化", "月別・年間・全期間レポート", "Cognito認証・本人別保存", "匿名比較・分類修正・学習"], 1.12, 3.48, 4.85, 1.9, size=15, gap=0.48)
    box(s, 6.82, 2.55, 5.7, 3.7, DARK, True)
    pill(s, "現在は停止", 7.16, 2.88, 1.42, fill="F4D7A4", color=DARK, size=12)
    bullet_list(s, ["公開URLと画面配信", "ログイン・API・クラウド保存", "実ユーザー平均との差の取得", "AWS上での再テスト"], 7.18, 3.48, 4.86, 1.9, size=15, color=WHITE, bullet_color="F4D7A4", gap=0.48)
    box(s, 0.78, 6.46, 11.74, 0.38, "F4D7A4", True)
    text(s, "残る主な課題：銀行CSV／返金・取消／二重計上の自動判定／削除API／最終取込警告の画面表示", 0.98, 6.52, 11.34, 0.23, size=10, color=DARK, align=PP_ALIGN.CENTER)
    footer(s, 11, 2)


def slide_12(prs):
    s = add_slide(prs, DARK)
    title_block(s, "11｜評価", "この展示で見ていただきたい4つの視点", "機能の多さではなく、理解しやすさ・安心感・行動へのつながりを評価してください。", dark=True)
    items = [
        ("01", "概要", "何を解決するものか分かった"),
        ("02", "使い方", "自分でもCSVを選べそうだと思った"),
        ("03", "安心感", "送る情報・送らない情報が伝わった"),
        ("04", "有用性", "振り返りが次の行動につながりそう"),
    ]
    for i, (num, heading, body) in enumerate(items):
        x = 0.78 + (i % 2) * 6.04
        y = 2.55 + (i // 2) * 1.38
        box(s, x, y, 5.66, 1.08, "243027", True, "48564D")
        text(s, num, x + 0.22, y + 0.24, 0.55, 0.3, size=12, color=LIME, bold=True)
        text(s, heading, x + 0.86, y + 0.18, 1.3, 0.36, size=16, color=WHITE, bold=True)
        text(s, body, x + 0.86, y + 0.56, 4.35, 0.3, size=12, color="D7DED9")
    box(s, 0.78, 5.62, 11.7, 0.78, LIME, True)
    rich_text(s, [("技術資料（別紙）を開く　", 16, True, DARK, TECH_DOC_LINK), ("セキュリティ・個人情報保護・技術選定・制約", 12, False, DARK, TECH_DOC_LINK)], 1.1, 5.8, 11.05, 0.4, valign=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.CENTER)
    text(s, "Google Drive取込後は、このボタンをGoogle DocsのURLへ差し替えます。", 1.2, 6.49, 10.9, 0.25, size=10, color="A8B2AC", align=PP_ALIGN.CENTER)
    footer(s, 12, 1.5, dark=True)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    prs.core_properties.title = "SpendOps Dashboard 展示資料"
    prs.core_properties.subject = "約20分の自由閲覧向け概要・使用方法・評価資料"
    prs.core_properties.author = "SpendOps Dashboard project"
    prs.core_properties.comments = "Google Slidesへのネイティブ取り込みを前提に作成"
    for slide_builder in [
        slide_01, slide_02, slide_03, slide_04, slide_05, slide_06,
        slide_07, slide_08, slide_09, slide_10, slide_11, slide_12,
    ]:
        slide_builder(prs)
    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(path)
